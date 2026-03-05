from pathlib import Path

from pptx import Presentation


def parse(path: Path, max_pages: int | None = None) -> str:
    """Parse PPTX files and extract text from slides."""
    prs = Presentation(str(path))
    parts: list[str] = []

    slides = list(prs.slides)
    if max_pages is not None:
        slides = slides[:max_pages]

    for i, slide in enumerate(slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append("\t".join(cells))
        if texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    if not parts:
        return "(No text content could be extracted from this PPTX)"

    return "\n\n".join(parts)
